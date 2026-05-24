import os
import logging
from typing import Optional, Dict, Any
from pydantic import BaseModel

# Optional library imports with robust fallback
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

logger = logging.getLogger("webrag.parser")

class ParsedDocument(BaseModel):
    filename: str
    text: str
    metadata: Dict[str, Any]
    page_count: int

class DocumentParser:
    @staticmethod
    def parse_file(file_path: str, filename: Optional[str] = None) -> ParsedDocument:
        """Parses various file types into clean text and metadata."""
        if not filename:
            filename = os.path.basename(file_path)
        
        ext = os.path.splitext(filename)[1].lower()
        logger.info(f"Parsing file {filename} (extension: {ext})")

        text = ""
        page_count = 1
        metadata = {
            "filename": filename,
            "extension": ext,
            "size": os.path.getsize(file_path)
        }

        try:
            if ext == ".pdf":
                text, page_count, pdf_meta = DocumentParser._parse_pdf(file_path)
                metadata.update(pdf_meta)
            elif ext in [".docx", ".doc"]:
                text, docx_meta = DocumentParser._parse_docx(file_path)
                metadata.update(docx_meta)
            elif ext in [".pptx", ".ppt"]:
                text, pptx_meta = DocumentParser._parse_pptx(file_path)
                metadata.update(pptx_meta)
            elif ext in [".xlsx", ".xls"]:
                text, xlsx_meta = DocumentParser._parse_xlsx(file_path)
                metadata.update(xlsx_meta)
            elif ext in [".html", ".htm"]:
                text = DocumentParser._parse_html(file_path)
            elif ext in [".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".ts", ".tsx", ".js", ".jsx", ".py", ".rs", ".go", ".c", ".cpp", ".h"]:
                text = DocumentParser._parse_text(file_path)
            else:
                # Default: try as text file
                try:
                    text = DocumentParser._parse_text(file_path)
                except Exception:
                    raise ValueError(f"Unsupported file format: {ext}")
            
            # Clean up the text a bit (remove excessive whitespace)
            lines = [line.strip() for line in text.splitlines()]
            text = "\n".join([line for line in lines if line])

        except Exception as e:
            logger.error(f"Error parsing file {filename}: {e}")
            raise e

        return ParsedDocument(
            filename=filename,
            text=text,
            metadata=metadata,
            page_count=page_count
        )

    @staticmethod
    def _parse_pdf(file_path: str) -> tuple[str, int, dict]:
        if not fitz:
            raise ImportError("PyMuPDF (fitz) is not installed.")
        
        doc = fitz.open(file_path)
        text_parts = []
        page_count = len(doc)
        
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            text_parts.append(page.get_text())
        
        full_text = "\n\n--- Page Break ---\n\n".join(text_parts)
        metadata = {
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "creator": doc.metadata.get("creator", ""),
            "producer": doc.metadata.get("producer", ""),
        }
        doc.close()
        return full_text, page_count, metadata

    @staticmethod
    def _parse_docx(file_path: str) -> tuple[str, dict]:
        if not docx:
            raise ImportError("python-docx is not installed.")
        
        doc = docx.Document(file_path)
        text_parts = []
        
        # Read paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                
        # Read tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                row_str = " | ".join([cell for cell in row_text if cell])
                if row_str:
                    text_parts.append(row_str)
                    
        full_text = "\n\n".join(text_parts)
        metadata = {
            "title": doc.core_properties.title or "",
            "author": doc.core_properties.author or "",
        }
        return full_text, metadata

    @staticmethod
    def _parse_pptx(file_path: str) -> tuple[str, dict]:
        if not pptx:
            raise ImportError("python-pptx is not installed.")
        
        prs = pptx.Presentation(file_path)
        text_parts = []
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
                
        full_text = "\n\n".join(text_parts)
        return full_text, {"slides_count": slide_count}

    @staticmethod
    def _parse_xlsx(file_path: str) -> tuple[str, dict]:
        if not openpyxl:
            raise ImportError("openpyxl is not installed.")
        
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell).strip() for cell in row if cell is not None]
                if row_vals:
                    sheet_text.append(" | ".join(row_vals))
            
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
                
        full_text = "\n\n".join(text_parts)
        return full_text, {"sheets_count": len(wb.sheetnames)}

    @staticmethod
    def _parse_html(file_path: str) -> str:
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 is not installed.")
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
            
        return soup.get_text()

    @staticmethod
    def _parse_text(file_path: str) -> str:
        # Try reading with utf-8, fallback to cp1252/latin-1
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        # Hard fallback with errors replaced
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
